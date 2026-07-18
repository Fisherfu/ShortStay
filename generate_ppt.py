import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_proposal_ppt():
    prs = Presentation()
    
    # 1. 設置幻燈片比例為 16:9 (13.33 x 7.5 英寸)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. 定義配色系統 (Premium Modern Palette)
    COLOR_BG_DARK = RGBColor(15, 23, 42)      # 深灰色 (Slate 900)
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # 淺灰色底 (Slate 50)
    COLOR_CARD = RGBColor(255, 255, 255)      # 卡片底色 (White)
    COLOR_PRIMARY = RGBColor(30, 41, 59)      # 主字體深藍 (Slate 800)
    COLOR_MUTED = RGBColor(100, 116, 139)     # 次字體灰色 (Slate 500)
    COLOR_TEAL = RGBColor(13, 148, 136)       # 主題綠 (Teal 600)
    COLOR_BLUE = RGBColor(37, 99, 235)        # 科技藍 (Blue 600)
    
    # 統一使用微軟正黑體
    FONT_NAME = "Microsoft JhengHei"
    
    # 輔助函式：設置背景色
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # 輔助函式：添加文字方塊
    def add_textbox(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.15):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.line_spacing = line_spacing
        return tf

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Background)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # 空白版面
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1, COLOR_BG_DARK)
    
    # 主標題
    add_textbox(slide1, Inches(1.5), Inches(2.2), Inches(10.33), Inches(1.2), 
                "客製化短租與房源管理平台", 
                font_size=44, font_color=RGBColor(255, 255, 255), bold=True)
    
    # 副標題
    add_textbox(slide1, Inches(1.5), Inches(3.4), Inches(10.33), Inches(0.8), 
                "專案開發與設計提案書\n以 Django + Python 實現高穩定、易擴充之解決方案", 
                font_size=20, font_color=COLOR_TEAL, bold=False, line_spacing=1.3)
    
    # 團隊資訊與日期
    add_textbox(slide1, Inches(1.5), Inches(5.2), Inches(5.0), Inches(0.5), 
                "提案團隊：Fisherfu 開發團隊\n日期：2026 年 7 月", 
                font_size=13, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)
    
    # ----------------------------------------------------
    # SLIDE 2: Background & Pain Points (Light Background)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_background(slide2, COLOR_BG_LIGHT)
    
    # 標題
    add_textbox(slide2, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "專案背景與市場痛點", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    # 3個痛點卡片設計
    card_width = Inches(3.64)
    card_height = Inches(4.5)
    card_top = Inches(1.6)
    
    pains = [
        ("痛點一：傳統套版缺乏彈性", "傳統套版網站（如 WordPress）的資料庫結構固定，難以應對複雜的自定義需求，例如「租客主動發布需求」與「房東留言即時媒合」的雙向互動邏輯。", COLOR_BLUE),
        ("痛點二：預訂流程不夠透明", "傳統訂房平台通常無即時狀態流轉。本專案提供自產預訂編號（BOOK+時間戳），支援「待確認、已確認、已取消、已完成」四種狀態，防堵訂單混亂。", COLOR_TEAL),
        ("痛點三：系統擴充與維護性差", "套版系統隨著外掛增多，速度變慢且有資安漏洞。我們採用 Django 獨立架構，高安全性防護（防 SQL Injection、XSS），且利於未來功能（如金流、API）擴充。", COLOR_PRIMARY)
    ]
    
    for i, (title, desc, accent_color) in enumerate(pains):
        left = Inches(0.8) + i * Inches(4.0)
        
        # 繪製卡片底色 (矩形)
        shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = RGBColor(226, 232, 240) # 邊框淺灰 (Slate 200)
        shape.line.width = Pt(1)
        
        # 卡片頂部裝飾條
        accent_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_width, Inches(0.15))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
        
        # 卡片文字內容
        # 標題
        add_textbox(slide2, left + Inches(0.3), card_top + Inches(0.4), card_width - Inches(0.6), Inches(0.6),
                    title, font_size=18, font_color=COLOR_PRIMARY, bold=True)
        # 描述
        add_textbox(slide2, left + Inches(0.3), card_top + Inches(1.2), card_width - Inches(0.6), Inches(3.0),
                    desc, font_size=14, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # ----------------------------------------------------
    # SLIDE 3: System Architecture (Light Background)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_background(slide3, COLOR_BG_LIGHT)
    
    add_textbox(slide3, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "技術選型與系統架構", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    # 左右雙欄設計：左欄架構說明，右欄運維特點
    # 左欄卡片
    left_x = Inches(0.8)
    col_width = Inches(5.6)
    
    shape_l = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(1.6), col_width, Inches(4.8))
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = COLOR_CARD
    shape_l.line.color.rgb = RGBColor(226, 232, 240)
    
    add_textbox(slide3, left_x + Inches(0.4), Inches(2.0), col_width - Inches(0.8), Inches(0.5),
                "🛠️ 高效後端與資料庫架構", font_size=18, font_color=COLOR_PRIMARY, bold=True)
    
    tech_text = (
        "• 後端技術 (Backend)：\n"
        "  採用企業級 Python 框架 Django 5.2。內建完整 ORM 模型與認證系統，縮短 50% 開發週期。\n\n"
        "• 資料庫 (Database)：\n"
        "  本地開發使用 SQLite 輕量資料庫加速敏捷開發；雲端生產環境採用 PostgreSQL 資料庫，支援高併發安全讀寫，防範數據遺失。"
    )
    add_textbox(slide3, left_x + Inches(0.4), Inches(2.7), col_width - Inches(0.8), Inches(3.3),
                tech_text, font_size=14, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # 右欄卡片
    right_x = Inches(6.9)
    shape_r = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, Inches(1.6), col_width, Inches(4.8))
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = COLOR_CARD
    shape_r.line.color.rgb = RGBColor(226, 232, 240)
    
    add_textbox(slide3, right_x + Inches(0.4), Inches(2.0), col_width - Inches(0.8), Inches(0.5),
                "🚀 現代化運維與自動部署", font_size=18, font_color=COLOR_PRIMARY, bold=True)
    
    ops_text = (
        "• 靜態檔案優化 (WhiteNoise)：\n"
        "  整合 WhiteNoise 靜態檔案引擎，不需複雜的伺服器 nginx 配置即可直接託管 CSS、JS、圖片，載入速度提升 200%。\n\n"
        "• 雲端自動部署 (Render + CI/CD)：\n"
        "  連結 GitHub 代碼倉庫。每次推送新功能，Render 自動執行 build.sh 進行安裝、遷移資料庫並重啟，實現不中斷自動部署。"
    )
    add_textbox(slide3, right_x + Inches(0.4), Inches(2.7), col_width - Inches(0.8), Inches(3.3),
                ops_text, font_size=14, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # ----------------------------------------------------
    # SLIDE 4: Core Features (Light Background)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_background(slide4, COLOR_BG_LIGHT)
    
    add_textbox(slide4, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "核心功能模組規劃", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    features = [
        ("01. 用戶權限系統", "自定義 CustomUser 角色系統，完美分離「房東、租客、管理員」，支持電話、頭像、註冊時間紀錄，具備健全的安全密碼校驗與重導向機制。"),
        ("02. 房源發布與管理", "支援房源詳細規格（臥室/浴室數/限住人數）、價格、地址與 GPS 經緯度，並採用靈活的 JSON 儲存硬體設施。提供封面圖設定與排序上傳。"),
        ("03. 智慧預訂與審核", "租客點選即自動計算住宿天數與總價，生成唯一預訂編號。房東專屬管理後台可一鍵審核「已確認、已取消」等預訂狀態，即時管控可用房源。"),
        ("04. 需求發布與留言", "租客可主動發布租房需求（如偏好城市、預算、期望格局與入住日期），房東或平台其他用戶可針對該需求進行專屬留言，達成即時雙向媒合。")
    ]
    
    f_width = Inches(5.6)
    f_height = Inches(2.1)
    
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        
        left = Inches(0.8) + col * Inches(6.1)
        top = Inches(1.6) + row * Inches(2.5)
        
        # 卡片底色
        shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, f_width, f_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        # 左側主題標記線
        marker = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), f_height)
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLOR_TEAL if col == 0 else COLOR_BLUE
        marker.line.fill.background()
        
        # 標題
        add_textbox(slide4, left + Inches(0.4), top + Inches(0.2), f_width - Inches(0.6), Inches(0.4),
                    title, font_size=16, font_color=COLOR_PRIMARY, bold=True)
        # 描述
        add_textbox(slide4, left + Inches(0.4), top + Inches(0.7), f_width - Inches(0.6), Inches(1.2),
                    desc, font_size=12, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # ----------------------------------------------------
    # SLIDE 5: Development Process (Light Background)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_background(slide5, COLOR_BG_LIGHT)
    
    add_textbox(slide5, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "專案開發時程與階段目標", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    steps = [
        ("第 1-2 週", "系統分析與介面設計", "• 需求細節確認與分析\n• 資料庫 Schema 詳細設計\n• UI/UX 視覺與網頁切版"),
        ("第 3-6 週", "核心功能開發與模組串接", "• 用戶認證與權限管理\n• 房源/圖片上傳與預訂模組\n• 租客需求留言版功能"),
        ("第 7-8 週", "系統測試與安全性優化", "• 單元測試與漏洞防範檢查\n• WhiteNoise 靜態檔案優化\n• 金流/地圖 API 模擬對接"),
        ("第 9 週", "雲端部署、驗收與交付", "• 生產環境 Render 部署與上線\n• 使用者驗收測試 (UAT)\n• 程式碼、文件轉移與教練培訓")
    ]
    
    step_width = Inches(2.7)
    step_height = Inches(4.5)
    step_top = Inches(1.8)
    
    for i, (time_lbl, title, bullets) in enumerate(steps):
        left = Inches(0.8) + i * Inches(2.9)
        
        # 繪製卡片
        shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, step_top, step_width, step_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        # 進度指示圈
        circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.1), step_top - Inches(0.3), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_TEAL if i < 3 else COLOR_BLUE
        circle.line.fill.background()
        
        # 時間標籤
        add_textbox(slide5, left, step_top + Inches(0.4), step_width, Inches(0.4),
                    time_lbl, font_size=16, font_color=COLOR_TEAL if i < 3 else COLOR_BLUE, bold=True, align=PP_ALIGN.CENTER)
        
        # 標題
        add_textbox(slide5, left + Inches(0.2), step_top + Inches(0.9), step_width - Inches(0.4), Inches(0.7),
                    title, font_size=14, font_color=COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        
        # 詳細點
        add_textbox(slide5, left + Inches(0.25), step_top + Inches(1.7), step_width - Inches(0.5), Inches(2.6),
                    bullets, font_size=11, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # ----------------------------------------------------
    # SLIDE 6: Maintenance & Warranty (Light Background)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_background(slide6, COLOR_BG_LIGHT)
    
    add_textbox(slide6, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "售後服務與保固維護", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    services = [
        ("🔧 12 個月免費系統保固", "自驗收日起提供 1 年免費程式 Bug 保固，保障平台穩定運作。保固範圍內程式故障皆由我方團隊第一時間修復。"),
        ("📊 完整錯誤日誌與錯誤捕獲", "Django 內建 Logging 機制自動紀錄重大錯誤，並可配置 Email/Webhook 自動警報，第一時間防堵程式崩潰，大幅降低停機損失。"),
        ("📖 健全的技術轉移與說明文件", "交付程式碼時附帶後台操作手冊與部署文檔。並提供 2 小時的線上教育訓練，確保您的團隊能輕鬆接手系統日常維護與內容管理。"),
        ("💾 資料庫定期雲端備份設計", "配合 Render 與雲端資料庫之備份策略，設計每日自動備份，防範因伺服器故障或惡意攻擊而導致的數據丟失，安全無虞。")
    ]
    
    s_width = Inches(5.6)
    s_height = Inches(2.1)
    
    for i, (title, desc) in enumerate(services):
        row = i // 2
        col = i % 2
        
        left = Inches(0.8) + col * Inches(6.1)
        top = Inches(1.6) + row * Inches(2.5)
        
        # 卡片底色
        shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, s_width, s_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        # 標題
        add_textbox(slide6, left + Inches(0.3), top + Inches(0.2), s_width - Inches(0.6), Inches(0.4),
                    title, font_size=16, font_color=COLOR_PRIMARY, bold=True)
        # 描述
        add_textbox(slide6, left + Inches(0.3), top + Inches(0.7), s_width - Inches(0.6), Inches(1.2),
                    desc, font_size=12, font_color=COLOR_MUTED, bold=False, line_spacing=1.3)

    # ----------------------------------------------------
    # SLIDE 7: Quotation & Timeline (Light Background)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_background(slide7, COLOR_BG_LIGHT)
    
    add_textbox(slide7, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.5), 
                "報價方案與付款時程", font_size=28, font_color=COLOR_PRIMARY, bold=True)
    
    # 左右兩個方案對比卡片
    p_width = Inches(5.6)
    p_height = Inches(4.8)
    p_top = Inches(1.6)
    
    # 方案 1: 原型快速上線版
    left_x = Inches(0.8)
    shape_p1 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, p_top, p_width, p_height)
    shape_p1.fill.solid()
    shape_p1.fill.fore_color.rgb = COLOR_CARD
    shape_p1.line.color.rgb = RGBColor(226, 232, 240)
    
    add_textbox(slide7, left_x + Inches(0.4), p_top + Inches(0.4), p_width - Inches(0.8), Inches(0.5),
                "方案 A：現有原型優化快速上線版", font_size=18, font_color=COLOR_TEAL, bold=True)
    
    add_textbox(slide7, left_x + Inches(0.4), p_top + Inches(0.9), p_width - Inches(0.8), Inches(0.5),
                "NT$ 60,000 ~ 80,000 (開發期：約 4 週)", font_size=16, font_color=COLOR_PRIMARY, bold=True)
    
    desc_p1 = (
        "• 使用現有 Django `ShortStay` 原型進行客製調整\n"
        "• 精美套版視覺調整（無從零客製設計師訪談）\n"
        "• 基本功能測試，並直接部署至 Render + PostgreSQL\n"
        "• 提供 3 個月系統 Bug 保固與 1 小時線上操作教學\n"
        "👉 適合需要短時間內（一個月內）上線 DEMO 的業主。"
    )
    add_textbox(slide7, left_x + Inches(0.4), p_top + Inches(1.6), p_width - Inches(0.8), Inches(3.0),
                desc_p1, font_size=12, font_color=COLOR_MUTED, bold=False, line_spacing=1.4)

    # 方案 2: 全功能高級客製化版
    right_x = Inches(6.9)
    shape_p2 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, p_top, p_width, p_height)
    shape_p2.fill.solid()
    shape_p2.fill.fore_color.rgb = COLOR_CARD
    shape_p2.line.color.rgb = RGBColor(226, 232, 240)
    
    add_textbox(slide7, right_x + Inches(0.4), p_top + Inches(0.4), p_width - Inches(0.8), Inches(0.5),
                "方案 B：全功能客製化設計版 (推薦)", font_size=18, font_color=COLOR_BLUE, bold=True)
    
    add_textbox(slide7, right_x + Inches(0.4), p_top + Inches(0.9), p_width - Inches(0.8), Inches(0.5),
                "NT$ 120,000 ~ 150,000 (開發期：約 9 週)", font_size=16, font_color=COLOR_PRIMARY, bold=True)
    
    desc_p2 = (
        "• 專屬 UI/UX 設計師參與訪談，打造獨一無二的前台視覺\n"
        "• 從零定製核心邏輯，支援第三方 Google/FB 登入\n"
        "• 串接第三方地圖 API (Google Maps) 及簡訊驗證 API\n"
        "• 完善的自動備份、伺服器監控與健全的技術交接文件\n"
        "• 提供 12 個月免費程式保固與 2 小時完整技術移轉培訓\n"
        "👉 適合需要建構品牌、高資安防護與長期營運的業主。"
    )
    add_textbox(slide7, right_x + Inches(0.4), p_top + Inches(1.6), p_width - Inches(0.8), Inches(3.0),
                desc_p2, font_size=12, font_color=COLOR_MUTED, bold=False, line_spacing=1.4)

    # 保存 PPTX 檔案
    ppt_path = "G:\\我的雲端硬碟\\OneDriveBackup\\交大資管所\\資料庫管理系統\\short_stay\\ShortStay_Development_Proposal.pptx"
    prs.save(ppt_path)
    print("PPTX generated successfully at: " + ppt_path)

if __name__ == '__main__':
    create_proposal_ppt()
